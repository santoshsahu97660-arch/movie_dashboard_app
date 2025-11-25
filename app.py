import streamlit as st
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt
import requests
import os

from io import BytesIO
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from pptx import Presentation
from pptx.util import Inches

# ---------- STREAMLIT CONFIG ----------
st.set_page_config(page_title="Movie Ratings Dashboard (Pro)", layout="wide")

st.title("🎬 Movie Ratings Dashboard (Pro)")
st.write("Upload Movie-Rating file and explore advanced analytics, downloads, and reports.")

# ---------- FILE UPLOAD ----------
uploaded_file = st.file_uploader("Movie-Rating wali CSV ya Excel file yaha upload karo", type=["csv", "xlsx"])

@st.cache_data
def load_data(file):
    if file.name.lower().endswith(".csv"):
        df = pd.read_csv(file)
    else:
        df = pd.read_excel(file)
    
    rename_map = {
        "Rotten Tomatoes Ratings %": "CriticRating",
        "Audience Ratings %": "AudienceRating",
        "Budget (million $)": "BudgetMillions",
        "Year of release": "Year"
    }
    df = df.rename(columns=rename_map)
    return df

def download_plot(fig):
    buf = BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    buf.seek(0)
    return buf

def generate_pdf_report(top3_highest, top3_lowest):
    buffer = BytesIO()
    c = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter

    c.setFont("Helvetica-Bold", 16)
    c.drawString(30, height - 30, "Movie Analytics Report")

    c.setFont("Helvetica", 12)
    c.drawString(30, height - 60, "📌 Top 3 Highest Rated Movies (Critic):")
    y = height - 80
    for name in top3_highest:
        c.drawString(40, y, name)
        y -= 15

    c.drawString(30, y - 10, "📌 Top 3 Lowest Rated Movies (Critic):")
    y -= 30
    for name in top3_lowest:
        c.drawString(40, y, name)
        y -= 15

    c.showPage()
    c.save()
    buffer.seek(0)
    return buffer

def generate_ppt_report(filtered, top_high_df, top_low_df, avg_genre_df):
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Movie Analytics Report"
    slide.placeholders[1].text = "Auto-generated from Streamlit dashboard"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Dataset Summary"
    tf = slide.placeholders[1].text_frame
    tf.text = f"Total filtered movies: {len(filtered)}"

    if "Genre" in filtered.columns:
        p = tf.add_paragraph()
        p.text = "Genres: " + ", ".join(sorted(map(str, filtered["Genre"].unique())))
        p.level = 1

    if "Year" in filtered.columns:
        p = tf.add_paragraph()
        p.text = "Years: " + ", ".join(sorted(map(str, filtered["Year"].unique())))
        p.level = 1

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Top 3 Highest & Lowest (Critic Rating)"
    tf = slide.placeholders[1].text_frame
    tf.text = "Top 3 Highest:"
    for _, row in top_high_df.iterrows():
        p = tf.add_paragraph()
        p.text = f"{row['Film']} – {row['CriticRating']}"
        p.level = 1

    p = tf.add_paragraph()
    p.text = "Top 3 Lowest:"
    for _, row in top_low_df.iterrows():
        q = tf.add_paragraph()
        q.text = f"{row['Film']} – {row['CriticRating']}"
        q.level = 1

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out

# -------------- MAIN APP --------------
if uploaded_file is None:
    st.info("📁 File upload karo (CSV/Excel) to start.")
else:
    movies = load_data(uploaded_file)

    st.subheader("📄 Data Preview")
    st.dataframe(movies.head())

    st.sidebar.header("🔍 Filters")
    genres = st.sidebar.multiselect("Genre", movies["Genre"].unique()) if "Genre" in movies.columns else []
    years = st.sidebar.multiselect("Year", movies["Year"].unique()) if "Year" in movies.columns else []

    filtered = movies[(movies["Genre"].isin(genres)) & (movies["Year"].isin(years))] if genres and years else movies.copy()

    st.write(f"🎬 Filtered Row Count: {filtered.shape[0]}")
    st.dataframe(filtered.head())

    csv_file = filtered.to_csv(index=False).encode("utf-8")
    st.download_button("⬇️ Download Filtered Data (CSV)", data=csv_file, file_name="filtered_movie_data.csv")

    if "CriticRating" in filtered.columns and not filtered.empty:
        top_3_highest = filtered.sort_values("CriticRating", ascending=False).head(3)[["Film", "CriticRating"]]
        top_3_lowest = filtered.sort_values("CriticRating", ascending=True).head(3)[["Film", "CriticRating"]]

        st.subheader("🏆 Top & Bottom 3 Movies")
        st.table(top_3_highest)
        st.table(top_3_lowest)

        st.download_button("📄 Download PDF Report", data=generate_pdf_report(top_3_highest["Film"].tolist(), top_3_lowest["Film"].tolist()), file_name="movie_report.pdf")

    # -------------- OMDb SECTION --------------
    st.markdown("---")
    st.subheader("🌐 OMDb / IMDb Style Live Data")

    api_key = st.sidebar.text_input("OMDb API Key (Hidden)", type="password", value=os.getenv("OMDB_API"))
    movie_query = st.sidebar.text_input("Movie Title")

    if st.sidebar.button("Fetch OMDb Data"):
        if not api_key or not movie_query:
            st.warning("API key aur movie title dono required hai.")
        else:
            try:
                resp = requests.get("http://www.omdbapi.com/", params={"t": movie_query, "apikey": api_key}, timeout=10)
                data = resp.json()
                if data.get("Response") == "True":
                    st.write(f"🎞 **{data.get('Title')} ({data.get('Year')})**")
                    st.write({
                        "IMDB Rating": data.get("imdbRating"),
                        "Genre": data.get("Genre"),
                        "Director": data.get("Director"),
                        "Actors": data.get("Actors"),
                        "Runtime": data.get("Runtime"),
                    })
                    st.write("📝 Plot:", data.get("Plot"))
                    if data.get("Poster") != "N/A":
                        st.image(data.get("Poster"), width=250)
                else:
                    st.warning("Movie not found.")
            except Exception as e:
                st.error(f"Error: {e}")
